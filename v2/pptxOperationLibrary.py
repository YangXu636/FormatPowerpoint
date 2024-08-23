import pptxOperationError as pptxOpError
from xyEnum import PpSaveAsFileType  # noqa
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu
import win32com.client
import os
from typing import Literal
import io
from PIL import Image
from betterTime import BetterTime as btime


class pptxOp:
    """Better Powerpoint Operation Library"""

    def __init__(self, file_path: str, t: int = 10) -> None:
        """
        __init__(file_path [,t=10])

        Initialize the pptxOp object.

        Parameters:
            file_path (str): The path of the Powerpoint file.
            t (int): The number of repeated trial and error attempts.

        Returns:
            None
        """
        self.file_path = file_path
        self.t = t

    def fileNew(self) -> None:
        """
        fileNew()

        Create a new Powerpoint file.

        Parameters:
            None

        Returns:
            None
        """
        t = self.t
        while t > 0:
            try:
                powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
                prs = powerpoint.Presentations.Add()
                if not os.path.exists(os.path.dirname(self.file_path) + "\\"):
                    os.makedirs(os.path.dirname(self.file_path) + "\\")
                prs.SaveAs(self.file_path)
                prs.Close()
                powerpoint.Quit()
                return
            except Exception as e:
                print(e)
                try:
                    prs.Close()
                except Exception:
                    pass
                powerpoint.Quit()
                t -= 1
        raise pptxOpError.FileCouldNotBeCreatedError(self.file_path)

    def fileConversion(self, new_file_format: str, new_file_path: str = None) -> None:
        """
        fileConversion(new_file_format,  [,new_file_path=None])

        Convert Powerpoint file to new_file_format.

        Parameters:
            new_file_format (str): The new file format. 'pptx', 'pdf', 'png', 'jpg', 'gif', 'tiff', 'bmp'
            new_file_path (str): The new file path. If None, the new file will be saved in the same directory as the original file with the same name as the original file but with the new file format.

        Returns:
            None
        """
        t = self.t
        if not os.path.exists(self.file_path):
            raise pptxOpError.FileNotFoundError(self.file_path)
        if new_file_path is None:
            new_file_path = os.path.splitext(self.file_path)[0] + "." + new_file_format
        else:
            new_file_path = (
                new_file_path + os.path.basename(self.file_path).split(".")[0]
            )
        print(new_file_path)
        if not os.path.exists(os.path.dirname(new_file_path)):
            os.makedirs(os.path.dirname(new_file_path))
        while t > 0:
            try:
                powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
                prs = powerpoint.Presentations.Open(self.file_path, WithWindow=False)
                # prs.SaveCopyAs(
                #     new_file_path, FileFormat=PpSaveAsFileType()[new_file_format]
                # )
                prs.SaveCopyAs(new_file_path + "." + new_file_format)
                prs.Close()
                powerpoint.Quit()
                return
            except Exception as e:
                print(e)
                try:
                    prs.Close()
                except Exception:
                    pass
                powerpoint.Quit()
                t -= 1
        raise pptxOpError.FileCouldNotBeConvertedError(self.file_path, new_file_format)

    def slideCopy(
        self,
        from_file_path: str,
        slide_num: int,
        new_slide_num: int = -1,
    ) -> None:
        """
        slideCopy(self, from_file_path, slide_num [,new_slide_num=len(self.slides)+1])

        Copy slide from \"from_file_path\" to \"self.file_path\".

        Parameters:
            from_file_path (str): The path of the Powerpoint file to copy from.
            slide_num (int): The number of the slide to copy.
            new_slide_num (int): The number of the new slide. If -1, the new slide will be added to the end of the Powerpoint file.

        Returns:
            None
        """
        t = self.t
        if not os.path.exists(from_file_path):
            raise pptxOpError.FileNotFoundError(from_file_path)
        if new_slide_num == -1:
            new_slide_num = len(self.file_path) + 1
        while t > 0:
            try:
                powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
                from_prs = powerpoint.Presentations.Open(
                    from_file_path, WithWindow=False
                )
                from_prs.Slides(slide_num).Copy()
                from_prs.Close()
                to_prs = powerpoint.Presentations.Open(self.file_path, WithWindow=False)
                to_prs.Slides.Paste(new_slide_num)
                to_prs.Save()
                to_prs.Close()
                return
            except Exception as e:
                error = e
                print(e)
                try:
                    from_prs.Close()
                except Exception:
                    pass
                try:
                    to_prs.Close()
                except Exception:
                    pass
                powerpoint.Quit()
                t -= 1
        raise pptxOpError.SlideCouldNotBeCopiedError(
            from_file_path, slide_num, self.file_path, new_slide_num, error
        )

    def slidesCount(self) -> int:
        """
        slidesCount() -> int

        Return the number of slides in the Powerpoint file.

        Parameters:
            None

        Returns:
            int: The number of slides in the Powerpoint file.
        """
        t = self.t
        while t > 0:
            try:
                powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
                prs = powerpoint.Presentations.Open(self.file_path, WithWindow=False)
                slides_count = len(prs.Slides)
                prs.Close()
                powerpoint.Quit()
                return slides_count
            except Exception as e:
                print(e)
                try:
                    prs.Close()
                except Exception:
                    pass
                powerpoint.Quit()
                t -= 1
        raise pptxOpError.FileCouldNotBeOpenedError(self.file_path)

    def slideSize(self) -> tuple:
        """
        slideSize() -> tuple(Pt, Pt)

        Return the size of the Powerpoint file in (width, height) format.

        Parameters:
            None

        Returns:
            tuple(Pt, Pt): The size of the Powerpoint file in (width, height) format.
        """
        t = self.t
        while t > 0:
            try:
                prs = pptx.Presentation(self.file_path)
                width, height = Emu(prs.slide_width).pt, Emu(prs.slide_height).pt
                del prs
                return (width, height)
            except Exception as e:
                print(e)
                try:
                    del prs
                except Exception:
                    pass
                t -= 1
        raise pptxOpError.FileCouldNotBeOpenedError(self.file_path)

    def slideTexts(self, slide_num: int) -> list[str]:
        """
        slideText(slide_num) -> list[str]

        Return the text of the slide in slide_num.

        Parameters:
            slide_num (int): The number of the slide to read.

        Returns:
            list[str]: The text of the slide in slide_num.
        """
        t = self.t
        while t > 0:
            try:
                powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
                prs = powerpoint.Presentations.Open(self.file_path, WithWindow=False)
                for i in prs.Slides(slide_num).Shapes:
                    if i.Type == MSO_SHAPE_TYPE.GROUP:
                        i.Ungroup()
                prs.Save()
                text = [
                    shape.TextFrame.TextRange.Text
                    for shape in prs.Slides(slide_num).Shapes
                    if shape.TextFrame.HasText
                ]
                prs.Close()
                powerpoint.Quit()
                return text
            except Exception as e:
                error = e
                print(e)
                try:
                    prs.Close()
                except Exception:
                    pass
                powerpoint.Quit()
                t -= 1
        raise pptxOpError.SlideCouldNotBeReadError(self.file_path, slide_num, error)

    def slidePictures(self, slide_num: int) -> list[bytes]:
        """
        slidePictures(slide_num) -> list[bytes]

        Return the blob of the picture in slide_num.

        Parameters:
            slide_num (int): The number of the slide to read.

        Returns:
            list: The blob of the picture in slide_num.
        """
        t = self.t
        while t > 0:
            try:
                powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
                prs = powerpoint.Presentations.Open(self.file_path, WithWindow=False)
                for i in prs.Slides(slide_num).Shapes:
                    if i.Type == MSO_SHAPE_TYPE.GROUP:
                        i.Ungroup()
                prs.Save()
                prs.Close()
                powerpoint.Quit()
                prs = pptx.Presentation(self.file_path)
                images = [
                    pic.image.blob
                    for pic in prs.slides[slide_num - 1].shapes
                    if pic.shape_type == MSO_SHAPE_TYPE.PICTURE
                ]
                del prs
                return images
            except Exception as e:
                error = e
                print(e)
                try:
                    del prs
                except Exception:
                    pass
                t -= 1
        raise pptxOpError.SlideCouldNotBeReadError(self.file_path, slide_num, error)

    def slide2Image(
        self,
        slide_num: int,
        image_path: str,
        image_name: str,
        image_format: Literal["png", "jpg", "gif", "tiff", "bmp"] = "png",
    ) -> None:
        """
        slide2Image(slide_num, image_path [,image_format="png"])

        Save the slide in slide_num as an image file.

        Parameters:
            slide_num (int): The number of the slide to save.
            image_path (str): The path of the image file to save.
            image_format (str): The format of the image file. 'png', 'jpg', 'gif', 'tiff', 'bmp'

        Returns:
            None
        """
        t = self.t
        while t > 0:
            try:
                if not os.path.exists(image_path):
                    os.makedirs(image_path)
                powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
                prs = powerpoint.Presentations.Open(self.file_path, WithWindow=False)
                for i in prs.Slides(slide_num).Shapes:
                    if i.Type == MSO_SHAPE_TYPE.GROUP:
                        i.Ungroup()
                prs.Save()
                prs.Slides(slide_num).Export(
                    image_path + f"\\{image_name}.{image_format}",
                    image_format.upper(),
                )
                prs.Close()
                powerpoint.Quit()
                return
            except Exception as e:
                error = e
                print(e)
                try:
                    prs.Close()
                except Exception:
                    pass
                powerpoint.Quit()
                t -= 1
        raise pptxOpError.SlideCouldNotBeSavedAsOtherFormatError(
            self.file_path, slide_num, image_format, error
        )

    def slide2Bytes(
        self,
        slide_num: int,
        image_format: Literal["png", "jpg", "gif", "tiff", "bmp"] = "png",
    ) -> bytes:
        """
        slide2Bytes(slide_num, image_format="png") -> bytes

        Return the slide in slide_num as a bytes object.

        Parameters:
            slide_num (int): The number of the slide to save.
            image_format (str): The format of the image file. 'png', 'jpg', 'gif', 'tiff', 'bmp'

        Returns:
            bytes: The slide in slide_num as a bytes object.
        """
        t = self.t
        while t > 0:
            try:
                tmpPath = os.getenv("TEMP")
                if not tmpPath:
                    tmpPath = f"C:\\Users\\{os.getlogin()}\\AppData\\Local\\Temp"
                powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
                prs = powerpoint.Presentations.Open(self.file_path, WithWindow=False)
                for i in prs.Slides(slide_num).Shapes:
                    if i.Type == MSO_SHAPE_TYPE.GROUP:
                        i.Ungroup()
                prs.Save()
                slide = prs.Slides(slide_num)
                slide.Export(
                    f"{tmpPath}\\SLIDE_{slide_num}.{image_format}", image_format.upper()
                )
                with open(f"{tmpPath}\\SLIDE_{slide_num}.{image_format}", "rb") as f:
                    data = f.read()
                os.remove(f"{tmpPath}\\SLIDE_{slide_num}.{image_format}")
                prs.Close()
                powerpoint.Quit()
                return io.BytesIO(data)
            except Exception as e:
                error = e
                print(e)
                try:
                    prs.Close()
                except Exception:
                    pass
                powerpoint.Quit()
                t -= 1
        raise pptxOpError.SlideCouldNotBeSavedAsOtherFormatError(
            self.file_path, slide_num, image_format, error
        )

    def textChange(self, slide_num: int, text: str, new_text: str) -> None:
        """
        textChange(slide_num, text, new_text) -> None

        Change the text of the slide in slide_num to new_text.



        Parameters:
            slide_num (int): The number of the slide to change.
            text (str): The text to change.
            new_text (str): The new text.

        Returns:
            None
        """
        t = self.t
        while t > 0:
            try:
                flag = False
                powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
                prs = powerpoint.Presentations.Open(self.file_path, WithWindow=False)
                for i in prs.Slides(slide_num).Shapes:
                    if i.Type == MSO_SHAPE_TYPE.GROUP:
                        i.Ungroup()
                for shape in prs.Slides(slide_num).Shapes:
                    if shape.HasTextFrame and shape.TextFrame.HasText:
                        if shape.TextFrame.TextRange.Text == text:
                            shape.TextFrame.TextRange.Text = new_text
                            flag = True
                if not flag:
                    raise pptxOpError.TextCouldNotBeChangedError(
                        self.file_path, slide_num, text, new_text, "Text not found"
                    )
                prs.Save()
                prs.Close()
                powerpoint.Quit()
                return
            except Exception as e:
                error = e
                print(e)
                try:
                    prs.Close()
                except Exception:
                    pass
                powerpoint.Quit()
                t -= 1
        raise pptxOpError.TextCouldNotBeChangedError(
            self.file_path, slide_num, text, new_text, error
        )

    def pictureAdd(
        self,
        slide_num: int,
        image_path: str | bytes,
        left: int,
        top: int,
        width: int = -1,
        height: int = -1,
    ) -> None:
        """
        pictureAdd(slide_num, image_path, left, top [,width=-1] [,height=-1]) -> None

        Add a picture to the slide in slide_num.

        Parameters:
            slide_num (int): The number of the slide to add the picture.
            image_path (str | bytes): The path of the image file or the bytes object of the image.
            left (int): The left position of the picture. The unit is in Pt.
            top (int): The top position of the picture. The unit is in Pt.
            width (int): The width of the picture. The unit is in Pt.
            height (int): The height of the picture. The unit is in Pt.

        Returns:
            None
        """
        t = self.t
        img = Image.open(image_path)
        if width <= 0:
            width = Emu(img.width).pt
        if height <= 0:
            height = Emu(img.height).pt
        img.close()
        while t > 0:
            try:
                powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
                prs = powerpoint.Presentations.Open(self.file_path, WithWindow=False)
                slide = prs.Slides(slide_num)
                slide.Shapes.AddPicture(
                    image_path,
                    left,
                    top,
                    width,
                    height,
                )
                prs.Save()
                prs.Close()
                powerpoint.Quit()
                return
            except Exception as e:
                error = e
                print(e)
                try:
                    prs.Close()
                except Exception:
                    pass
                powerpoint.Quit()
                t -= 1
        raise pptxOpError.PictureCouldNotBeAddedError(
            self.file_path, slide_num, image_path, left, top, error
        )

    def pictureChange(self, slide_num: int, old_image_bytes, new_image_bytes) -> None:
        """
        pictureChange(slide_num, old_image_bytes, new_image_bytes) -> None

        Change the picture of the slide in slide_num to new_image_bytes.

        Parameters:
            slide_num (int): The number of the slide to change.
            old_image_bytes (blob): The bytes object of the old image.
            new_image_bytes (blob): The bytes object of the new image.

        Returns:
            None
        """
        t = self.t
        while t > 0:
            try:
                prs = pptx.Presentation(self.file_path)
                prs_size = (Emu(prs.slide_width).pt, Emu(prs.slide_height).pt)
                del prs
                flag = False
                powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
                prs = powerpoint.Presentations.Open(self.file_path, WithWindow=False)
                for i in prs.Slides(slide_num).Shapes:
                    if i.Type == MSO_SHAPE_TYPE.GROUP:
                        i.Ungroup()
                prs.Save()
                prs.Close()
                powerpoint.Quit()
                prs = pptx.Presentation(self.file_path)
                for i in range(len(prs.slides[slide_num - 1].shapes)):
                    shape = prs.slides[slide_num - 1].shapes[i]
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        if shape.image.blob == old_image_bytes:
                            flag = True
                            old_image_size = i + 1
                if not flag:
                    raise pptxOpError.PictureCouldNotBeChangedError(
                        self.file_path,
                        slide_num,
                        old_image_bytes,
                        new_image_bytes,
                        "Picture not found",
                    )
                del prs
                powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
                prs = powerpoint.Presentations.Open(self.file_path, WithWindow=False)
                slide = prs.Slides(slide_num)
                shape = slide.Shapes(old_image_size)
                old_image_size = [
                    shape.Width,
                    shape.Height,
                    shape.Left,
                    shape.Top,
                ]
                shape.Delete()
                if not isinstance(new_image_bytes, str):
                    new_image = new_image_bytes
                    new_image = Image.open(new_image)
                    new_image_bytes = (
                        os.getenv("TEMP") + f"\\new_image_{btime.timeName()}.png"
                    )
                    new_image.save(new_image_bytes)
                    new_image.close()
                new_image = Image.open(new_image_bytes)
                new_image_size = new_image.size
                new_image.close()
                new_image_size = (
                    new_image_size[0] / new_image_size[1] * old_image_size[1],
                    old_image_size[1],
                )
                if new_image_size[0] > prs_size[0]:
                    new_image_size = (
                        prs_size[0],
                        prs_size[0] / new_image_size[0] * new_image_size[1],
                    )
                slide.Shapes.AddPicture(
                    new_image_bytes,
                    True,
                    True,
                    (prs_size[0] - new_image_size[0]) / 2.0000,
                    old_image_size[3]
                    + (old_image_size[1] - new_image_size[1]) / 2.0000,
                    new_image_size[0],
                    new_image_size[1],
                )
                prs.Save()
                prs.Close()
                powerpoint.Quit()
                return
            except Exception as e:
                error = e
                print(e)
                try:
                    prs.Close()
                except Exception:
                    try:
                        del prs
                    except Exception:
                        pass
                try:
                    powerpoint.Quit()
                except Exception:
                    pass
                t -= 1
        raise pptxOpError.PictureCouldNotBeChangedError(
            self.file_path, slide_num, old_image_bytes, new_image_bytes, error
        )

    def componentDeleteText(self, slide_num: int, component_content: str) -> None:
        """
        componentDelete(slide_num, component_content) -> None

        Delete the component in slide_num.

        Parameters:
            slide_num (int): The number of the slide to delete.
            component_content (str): The content of the component to delete.

        Returns:
            None
        """
        t = self.t
        while t > 0:
            try:
                powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
                prs = powerpoint.Presentations.Open(self.file_path, WithWindow=False)
                for i in prs.Slides(slide_num).Shapes:
                    if i.Type == MSO_SHAPE_TYPE.GROUP:
                        i.Ungroup()
                prs.Save()
                delete_shape = None
                for shape in prs.Slides(slide_num).Shapes:
                    if shape.TextFrame.HasText:
                        if shape.TextFrame.TextRange.Text == component_content:
                            delete_shape = shape
                            break
                if delete_shape is None:
                    raise pptxOpError.ComponentCouldNotBeDeletedError(
                        self.file_path,
                        slide_num,
                        component_content,
                        "Component not found",
                    )
                try:
                    delete_shape.Delete()
                except Exception as e:
                    raise pptxOpError.ComponentCouldNotBeDeletedError(
                        self.file_path,
                        slide_num,
                        component_content,
                        f"Component could not be deleted due to error: {e}",
                    )
                prs.Save()
                prs.Close()
                powerpoint.Quit()
                return
            except Exception as e:
                error = e
                print(e)
                try:
                    prs.Close()
                except Exception:
                    try:
                        del prs
                    except Exception:
                        pass
                try:
                    powerpoint.Quit()
                except Exception:
                    pass
                t -= 1
        raise pptxOpError.ComponentCouldNotBeDeletedError(
            self.file_path, slide_num, component_content, error
        )

    def componentDeletePicture(self, slide_num: int, component_content: bytes) -> None:
        """
        componentDeletePicture(slide_num, component_content) -> None

        Delete the picture in slide_num.

        Parameters:
            slide_num (int): The number of the slide to delete.
            component_content (bytes): The bytes object of the picture to delete.

        Returns:
            None
        """
        t = self.t
        while t > 0:
            try:
                powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
                prs = powerpoint.Presentations.Open(self.file_path, WithWindow=False)
                for i in prs.Slides(slide_num).Shapes:
                    if i.Type == MSO_SHAPE_TYPE.GROUP:
                        i.Ungroup()
                prs.Save()
                prs.Save()
                prs.Close()
                ppt = pptx.Presentation(self.file_path)
                delete_shape_num = None
                for i in range(len(ppt.slides[slide_num - 1].shapes)):
                    shape = ppt.slides[slide_num - 1].shapes[i]
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        if shape.image.blob == component_content:
                            delete_shape_num = i + 1
                            break
                del ppt
                if delete_shape_num is None:
                    raise pptxOpError.PictureCouldNotBeDeletedError(
                        self.file_path,
                        slide_num,
                        component_content,
                        "Picture not found",
                    )
                try:
                    prs = powerpoint.Presentations.Open(
                        self.file_path, WithWindow=False
                    )
                    prs.Slides(slide_num).Shapes(delete_shape_num).Delete()
                except Exception as e:
                    raise pptxOpError.ComponentCouldNotBeDeletedError(
                        self.file_path,
                        slide_num,
                        component_content,
                        f"Picture could not be deleted due to error: {e}",
                    )
                prs.Save()
                prs.Close()
                powerpoint.Quit()
                return
            except Exception as e:
                error = e
                print(e)
                try:
                    prs.Close()
                except Exception:
                    try:
                        del prs
                    except Exception:
                        pass
                try:
                    powerpoint.Quit()
                except Exception:
                    pass
                t -= 1
        raise pptxOpError.ComponentCouldNotBeDeletedError(
            self.file_path, slide_num, component_content, error
        )

    def componentDelete(
        self,
        slide_num: int,
        component_content: str | bytes | list | tuple,
        reportErrors: bool = True,
    ) -> None:
        """
        componentDelete(slide_num, component_content[, reportErrors=True]) -> None

        Delete the component(s) in slide_num.

        Parameters:
            slide_num (int): The number of the slide to delete.
            component_content (str | bytes | list): The content(s) of the component to delete.
            reportErrors (bool): Whether to report errors or not.

        Returns:
            None
        """
        if reportErrors:
            if isinstance(component_content, str):
                self.componentDeleteText(slide_num, component_content)
            elif isinstance(component_content, bytes):
                self.componentDeletePicture(slide_num, component_content)
            elif isinstance(component_content, list):
                for content in component_content:
                    self.componentDelete(slide_num, content, reportErrors)
            elif isinstance(component_content, tuple):
                for content in component_content:
                    self.componentDelete(slide_num, content, reportErrors)
            else:
                raise TypeError("component_content must be str, bytes, list or tuple")
        else:
            try:
                if isinstance(component_content, str):
                    self.componentDeleteText(slide_num, component_content)
                elif isinstance(component_content, bytes):
                    self.componentDeletePicture(slide_num, component_content)
                elif isinstance(component_content, list):
                    for content in component_content:
                        self.componentDelete(slide_num, content, reportErrors)
            except Exception as e:  # noqa
                pass
