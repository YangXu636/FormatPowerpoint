import shutil
import zipfile
import os
import win32com.client
import win32clipboard
import pywintypes
import time
from pptx import Presentation
from pptx.dml.color import _NoneColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Pt
import re
from PIL import Image


def copy_file(src, dst):
    shutil.copy(src, dst)


def remove_file(path):
    if os.path.isdir(path):
        shutil.rmtree(path)
        os.mkdir(path)
    else:
        os.remove(path)


def ZipExtract(zip_file, dir):
    if not os.path.exists(dir):
        os.mkdir(dir)
    with zipfile.ZipFile(zip_file, "r") as zip_ref:
        zip_ref.extractall(path=dir)


def ZipCompress(zip_file, dir):
    if os.path.exists(zip_file):
        os.remove(zip_file)
    zip = zipfile.ZipFile(zip_file, "w", zipfile.ZIP_DEFLATED)
    for path, dirnames, filenames in os.walk(dir):  # type:ignore
        # 去掉目标跟路径，只对目标文件夹下边的文件及文件夹进行压缩
        fpath = path.replace(dir, "")
        for filename in filenames:
            zip.write(os.path.join(path, filename), os.path.join(fpath, filename))
    zip.close()


def format_ppt(
    mb_file,
    need_format_file,
    target_file,
    start_index,
    finish_index,
    log,
    success_log,
    fail_log,
):
    log(
        f"format_ppt start\n{mb_file = }\n{need_format_file = }\n{log = }\n{success_log = }\n{fail_log = }"
    )
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1
    prs_size = (0, 0)
    try:
        prs_size = set_ppt_title(mb_file, need_format_file, target_file, log)
        print(f"prs_size = {prs_size}")
        success_log("title")
    except Exception as e:
        fail_log(f"title, {e =}")
    mbChapterIndex = get_ppt_chapter_slideIndex(mb_file, log)
    log(f"{mbChapterIndex = }")
    mbCatalogueIndex = get_ppt_catalogues_slideIndex(mb_file, log)
    log(f"{mbCatalogueIndex = }")
    # mbOnlyImageIndex = get_ppt_only_image_slideIndex(mb_file, log)
    # log(f"{mbOnlyImageIndex = }")
    # mbOnlyParagraphIndex = get_ppt_only_paragraph_slideIndex(mb_file, log)
    # log(f"{mbOnlyParagraphIndex = }")
    # mbImageAndParagraphIndex = get_ppt_image_paragraph_slideIndex(mb_file, log)
    # log(f"{mbImageAndParagraphIndex = }")

    mbOnlyImageIndex = 5
    mbOnlyParagraphIndex = 6
    mbImageAndParagraphIndex = 4
    mbImageAndExplainIndex = 7

    if mbChapterIndex == -1:
        log("can't find chapter slide")
    if mbCatalogueIndex == -1:
        log("can't find catalogue slide")
    if mbOnlyImageIndex == -1:
        log("can't find only image slide")
    if mbOnlyParagraphIndex == -1:
        log("can't find only paragraph slide")
    if mbImageAndParagraphIndex == -1:
        log("can't find image and paragraph slide")
    if mbImageAndExplainIndex == -1:
        log("can't find image and explain slide")
    prs = Presentation(need_format_file)
    text = ""
    image_count = []
    have_paragraph = []
    NowCatalogue = ""
    imageNum = 1
    # for i in range(1, get_needFormatPpt_count(need_format_file, log)):
    for i in range(
        1,
        min(finish_index, get_needFormatPpt_count(need_format_file, log)),
    ):
        print(f"1 and {i = }")
        flag_addImage = False
        try:
            prses = powerpoint.Presentations.Open(need_format_file)
        except pywintypes.com_error as e:
            print(f"open ppt {i} error!, {e = }")
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = 1
            prses = powerpoint.Presentations.Open(need_format_file)
        except Exception as e:
            print(f"open ppt {i}, {e = }")
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = 1
            prses = powerpoint.Presentations.Open(need_format_file)
        for j in range(len(prs.slides[i].shapes)):
            if prs.slides[i].shapes[j].has_text_frame:
                text += "#fgf#" + prs.slides[i].shapes[j].text_frame.text
        for shape in prses.Slides(i + 1).Shapes:
            if (
                shape.TextFrame.HasText
                and shape.TextFrame.TextRange.ParagraphFormat.Bullet.Visible
            ):
                have_paragraph.append(shape.TextFrame.TextRange.Text)
            if shape.Type == MSO_SHAPE_TYPE.PICTURE and not flag_addImage:
                image_count.append(f"image{imageNum}.jpeg")
                flag_addImage = True
                imageNum += 1
        flag_addImage = False
        prses.Close()
        log(f"第{i}张ppt的文本内容为 {text}")
        log(f"第{i}张ppt的图片数量为 {image_count}")
        log(f"第{i}张ppt段落 {have_paragraph}")
        print(f"text = {text}")
        if i < start_index:
            text = ""
            image_count = []
            have_paragraph = []
            NowCatalogue = ""
            continue
        if "第" in text and "章" in text:
            # try:
            #     set_ppt_chapter(
            #         mb_file, need_format_file, target_file, mbChapterIndex, i, log
            #     )
            #     success_log(f"chapter {i}")
            # except Exception as e:
            #     fail_log(f"chapter {i}  {e = }")
            try:
                textList = text.split("#fgf#")
                for texts in textList:
                    if (
                        "第" in texts
                        and "章" in texts
                        and is_number(
                            re.findall(r"第.*?章", texts)[0][1:-1]
                        )  # r"(第[\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10}(章|节))"
                        and len(texts) <= 20
                    ):
                        set_ppt_chapter(
                            mb_file,
                            need_format_file,
                            target_file,
                            mbChapterIndex,
                            texts,
                            log,
                        )
                        success_log(f"chapter {i}")
                        NowCatalogue = re.findall(r"第.*?章", texts)[0][1:-1]
                        break
                if re.findall(
                    r"(第[\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10}节)", text, re.U
                ) and re.findall(
                    r"((第[\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10}章 ).*?(?=#fgf#))",
                    text,
                    re.U,
                ):
                    set_ppt_catalogue(
                        mb_file,
                        target_file,
                        mbCatalogueIndex,
                        re.search(
                            r"((第[\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10}章 ).*?(?=#fgf#))",  # ((?<=第[\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10}章 ).*?(?=#fgf#))
                            text,
                        )
                        .group()
                        .split(" ")[
                            -1
                        ],  # (第[\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10}章.*(?=#fgf#))
                        "\n".join(
                            [
                                j
                                for i in text.split("#fgf#")
                                if "第" in i and "节" in i
                                for j in i.split("\n")
                            ]
                        ),
                        log,
                    )
                    success_log(f"catalogueOne {i}")
                    NowCatalogue = (
                        re.search(
                            r"((第[\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10}章 ).*?(?=#fgf#))",  # ((?<=第[\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10}章 ).*?(?=#fgf#))
                            text,
                        )
                        .group()
                        .split(" ")[-1]
                    )
            except Exception as e:
                fail_log(f"chapterOne {i} & {e = }")
        elif "第" in text and "节" in text:
            try:
                set_ppt_catalogue(
                    mb_file,
                    target_file,
                    mbCatalogueIndex,
                    re.search(
                        r"((第[\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10}节 ).*?(?=#fgf#))",  # ((?<=第[\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10}章 ).*?(?=#fgf#))
                        text,
                    )
                    .group()
                    .split(" ")[
                        -1
                    ],  # (第[\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10}章.*(?=#fgf#))
                    "\n".join(
                        [
                            j
                            for i in text.split("#fgf#")
                            if "、" in i
                            for j in i.split("\n")
                            if "、" in j
                        ]
                    ),
                    log,
                )
                success_log(f"catalogueTwo {i}")
                NowCatalogue = (
                    re.search(
                        r"((第[\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10}节 ).*?(?=#fgf#))",  # ((?<=第[\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10}章 ).*?(?=#fgf#))
                        text,
                    )
                    .group()
                    .split(" ")[-1]
                )
            except Exception as e:
                fail_log(f"catalogueTwo {i}  {e = }")
        elif (
            re.findall(r"([\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10})、.*?(?=#fgf#)", text)
            and re.findall(
                r"[\(（]([\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10})[\)）]、.*?(?=\n|$)",
                text,
            )  # re.finditer(r"（([\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10})）、.*?(?=\s|$)", text)
        ):
            try:
                set_ppt_catalogue(
                    mb_file,
                    target_file,
                    mbCatalogueIndex,
                    re.search(
                        r"([\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10})、.*?(?=#fgf#)",
                        text,
                    )
                    .group()
                    .split("、")[-1],
                    "\n".join(
                        [
                            i.group()
                            for i in re.finditer(
                                r"[\(（]([\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10})[\)）]、.*?(?=\n|$)",
                                text,
                            )
                        ]
                    ),
                    log,
                )
                success_log(f"catalogueThree {i}")
                NowCatalogue = (
                    re.search(
                        r"([\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10})、.*?(?=#fgf#)",
                        text,
                    )
                    .group()
                    .split("、")[-1]
                )
            except Exception as e:
                fail_log(f"catalogueThree {i}  {e = }")
        elif re.findall(
            r"[\(（]([\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10})[\)）]、.*?(?=#fgf#)",
            text,
        ) and re.findall(r"([1-9]{1,10})..*?(?=\n|$)", text):
            try:
                set_ppt_catalogue(
                    mb_file,
                    target_file,
                    mbCatalogueIndex,
                    re.search(
                        r"[\(（]([\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10})[\)）]、.*?(?=#fgf#)",
                        text,
                    )
                    .group()
                    .split("、")[-1],
                    "\n".join(
                        [
                            i.group()
                            for i in re.finditer(
                                r"([1-9]{1,10})..*?(?=\n|$)",
                                text,
                            )
                        ]
                    ),
                    log,
                )
                success_log(f"catalogueFour {i}")
                NowCatalogue = (
                    re.search(
                        r"[\(（]([\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10})[\)）]、.*?(?=#fgf#)",
                        text,
                    )
                    .group()
                    .split("、")[-1]
                )
            except Exception as e:
                fail_log(f"catalogueFour {i}  {e = }")
        elif re.findall(
            r"(\(|（)见{0,}(图|表)\d{1,10}-\d{1,100}[A-Z | a-z]{0,}(）|\))", text
        ) and (
            text.count("#fgf#") == 1
            or len([ik for ik in text.split("#fgf#") if ik]) == 1
        ):
            try:
                set_ppt_only_image(
                    mb_file,
                    need_format_file,
                    target_file,
                    mbOnlyImageIndex,
                    NowCatalogue,
                    image_count,
                    re.search(
                        r"(?<=#fgf#).*?(?:(\(|（)见{0,}(图|表)\d{1,10}-\d{1,100}[A-Z | a-z]{0,}(）|\)))",
                        text,
                    ).group(),
                    prs_size,
                    log,
                )
                success_log(f"only image {i}")
            except Exception as e:
                fail_log(f"only image {i}  {e = }")
        elif image_count and text.count("#fgf#") >= 2:
            try:
                set_ppt_image_explain(
                    mb_file,
                    need_format_file,
                    target_file,
                    mbImageAndExplainIndex,
                    NowCatalogue,
                    image_count,
                    re.search(
                        r"(?<=#fgf#).*?(?:(\(|（)见{0,}(图|表)\d{1,10}-\d{1,100}[A-Z | a-z]{0,}(）|\)))",
                        text,
                    ).group(),
                    have_paragraph[0]
                    if have_paragraph
                    else [i for i in text.strip().split("#fgf#") if i][1],
                    prs_size,
                    log,
                )
                success_log(f"image and explain {i}")
            except Exception as e:
                fail_log(f"image and explain {i}  {e = }")
        elif have_paragraph:
            print(f"have_paragraph = {have_paragraph}")
            try:
                set_ppt_only_paragraph(
                    mb_file, target_file, mbOnlyParagraphIndex, text, NowCatalogue, log
                )
                success_log(f"only paragraph One {i}")
            except Exception as e:
                fail_log(f"only paragraph One {i}  {e = }")
        elif text.find("#fgf#") == 1:
            try:
                set_ppt_only_paragraph(
                    mb_file, target_file, mbOnlyParagraphIndex, text, NowCatalogue, log
                )
                success_log(f"only paragraph Two {i}")
            except Exception as e:
                fail_log(f"only paragraph Two {i}  {e = }")
        else:
            fail_log(f"image or paragraph {i}")
            pass
        text = ""
        image_count = []
        have_paragraph = []
        if i % 5 == 0:
            time.sleep(1)
        if i % 100 == 0:
            time.sleep(5)
    del prs
    log(f"format_ppt end, life = {target_file}\n\n\n")
    return


def set_ppt_only_image(
    mb_file,
    nf_file,
    target_file,
    mbIndex,
    catalogue,
    images,
    image_jianjie,
    prs_size,
    log,
):
    print(f"{images = }")
    time.sleep(1)
    powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
    powerpoint.Visible = 1
    while 10:
        try:
            mb = powerpoint.Presentations.Open(mb_file)
            mb.Slides(mbIndex).Copy()
            print("Copy!")
            mb.Close()
            print("mbClose!")
            index = get_needFormatPpt_count(target_file, log) + 1
            target = powerpoint.Presentations.Open(target_file)
            print("OpenTarget!")
            target.Slides.Paste(index)
            print("Paste!")
            raise RuntimeError("break")
        except RuntimeError:
            break
        except Exception:
            print("Copy or Paste Error, Try again!")
            try:
                target.Close()
            except Exception:
                pass
            powerpoint.Quit()
            powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
    picHeight = 100
    introduction = None
    picNeedDel = None
    while 10:
        try:
            print(f"{target.Slides(index).Shapes.Count = }")
            for j in target.Slides(index).Shapes:
                print(f"{j.Type = } {j.Name = }")
                if j.Type == MSO_SHAPE_TYPE.GROUP:
                    j.Ungroup()
                if j.Type == MSO_SHAPE_TYPE.PICTURE and j.Name != "Picture 3":
                    picHeight = j.Height
                    if not picNeedDel:
                        picNeedDel = j
                if j.TextFrame.HasText:
                    if "SubTitle" in j.TextFrame.TextRange.Text:
                        j.TextFrame.TextRange.Text = catalogue
                    elif "Introduction" in j.TextFrame.TextRange.Text:
                        introduction = j
                        j.TextFrame.TextRange.Text = image_jianjie
            raise RuntimeError("break")
        except RuntimeError:
            break
        except Exception as e:
            print(f"set_ppt_only_image error {e = }, Try again!")
            picHeight = 100
            introduction = None
            picNeedDel = None
            try:
                target.Close()
            except Exception:
                pass
            powerpoint.Quit()
            powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
            target = powerpoint.Presentations.Open(target_file)
    if picNeedDel:
        picNeedDel.Delete()
    # picPath = [nf_file[:-5] + "\\ppt\\media\\" + images[i] for i in range(len(images))]
    # picPosition = get_picture_position(picPath, prs_size, picHeight, log)
    # print(f"{picPosition = }")
    # for i in picPosition:
    #     target.Slides(index).Shapes.AddPicture(i[0], True, True, i[1], i[2], i[3], i[4])
    picPath = nf_file[:-5] + "\\ppt\\media\\" + images[0]
    picPos = get_picture_position(picPath, prs_size, picHeight, log)
    print(f"{picPos = }")
    target.Slides(index).Shapes.AddPicture(
        picPos[0], True, True, picPos[1], picPos[2], picPos[3], picPos[4]
    )
    introduction.Left = (Emu(prs_size[0]).pt - introduction.Width) / 2.0000
    target.Save()
    print("SaveTarget!")
    target.Close()
    print("CloseTarget!")
    return


def set_ppt_only_paragraph(
    mb_file,
    target_file,
    mbIndex,
    paragraphs,
    NowCatalogue,
    log,
):
    print(f"{paragraphs = }")
    time.sleep(1)
    powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
    powerpoint.Visible = 1
    while 10:
        try:
            mb = powerpoint.Presentations.Open(mb_file)
            mb.Slides(mbIndex).Copy()
            print("Copy!")
            mb.Close()
            print("mbClose!")
            index = get_needFormatPpt_count(target_file, log) + 1
            target = powerpoint.Presentations.Open(target_file)
            print("OpenTarget!")
            target.Slides.Paste(index)
            print("Paste!")
            raise RuntimeError("break")
        except RuntimeError:
            break
        except Exception:
            print("Copy or Paste Error, Try again!")
            try:
                target.Close()
            except Exception:
                pass
            powerpoint.Quit()
            powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
    texts = [i for i in paragraphs.rstrip().split("#fgf#") if i]
    catalogue = ""
    paragraph = ""
    if re.findall(r"([\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10})、.*?(?=#fgf#)", texts[0]):
        catalogue = re.search(
            r"([\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10})、.*?(?=#fgf#)", texts[0]
        ).group()
    elif re.findall(
        r"[\(（]([\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10})[\)）]、.*?(?=#fgf#)", texts[0]
    ):
        catalogue = re.search(
            r"[\(（]([\u4e00-\u9fa5\u767e\u5343\u96f6]{1,10})[\)）]、.*?(?=#fgf#)",
            texts[0],
        ).group()
    elif re.findall(r"([1-9]{1,10})..*?(?=#fgf#)", texts[0]):
        catalogue = re.search(r"([1-9]{1,10})..*?(?=#fgf#)", texts[0]).group()
    elif "本章概述" in texts[0]:
        catalogue = "本章概述"
    else:
        catalogue = texts[0]
    if len(texts) > 1:
        paragraph = "\n".join(texts[1:])
    if len(catalogue) > len(paragraph) == 0:
        paragraph = catalogue
        catalogue = NowCatalogue
    while 10:
        try:
            print(f"{target.Slides(index).Shapes.Count = }")
            for j in target.Slides(index).Shapes:
                print(f"{j.Type = } {j.Name = }")
                if j.Type == MSO_SHAPE_TYPE.GROUP:
                    j.Ungroup()
                if j.TextFrame.HasText:
                    if "SubTitle" in j.TextFrame.TextRange.Text:
                        j.TextFrame.TextRange.Text = catalogue
                    elif "aragraph" in j.TextFrame.TextRange.Text:
                        j.TextFrame.TextRange.Text = paragraph
            raise RuntimeError("break")
        except RuntimeError:
            break
        except Exception as e:
            print(f"set_ppt_only_paragraph error {e = }, Try again!")
            catalogue = ""
            paragraph = ""
            try:
                target.Close()
            except Exception:
                pass
            powerpoint.Quit()
            powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
            target = powerpoint.Presentations.Open(target_file)
    target.Save()
    print("SaveTarget!")
    target.Close()
    print("CloseTarget!")
    return


def set_ppt_image_paragraph(
    mb_file, target_file, mbIndex, catalogue, images, paragraphs, log
):
    print(f"{images = } {paragraphs = }")
    pass


def set_ppt_image_explain(
    mb_file,
    nf_file,
    target_file,
    mbIndex,
    catalogue,
    images,
    image_jianjie,
    explain,
    prs_size,
    log,
):
    print(f"{images = } {explain = }")
    time.sleep(1)
    powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
    powerpoint.Visible = 1
    while 10:
        try:
            mb = powerpoint.Presentations.Open(mb_file)
            mb.Slides(mbIndex).Copy()
            print("Copy!")
            mb.Close()
            print("mbClose!")
            index = get_needFormatPpt_count(target_file, log) + 1
            target = powerpoint.Presentations.Open(target_file)
            print("OpenTarget!")
            target.Slides.Paste(index)
            print("Paste!")
            raise RuntimeError("break")
        except RuntimeError:
            break
        except Exception:
            print("Copy or Paste Error, Try again!")
            try:
                target.Close()
            except Exception:
                pass
            powerpoint.Quit()
            powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
    picHeight = 100
    introduction = None
    picNeedDel = None
    picPosition = ()
    Explain = None
    while 10:
        try:
            print(f"{target.Slides(index).Shapes.Count = }")
            for j in target.Slides(index).Shapes:
                print(f"{j.Type = } {j.Name = }")
                if j.Type == MSO_SHAPE_TYPE.GROUP:
                    j.Ungroup()
                if j.Type == MSO_SHAPE_TYPE.PICTURE and j.Name != "Picture 3":
                    picHeight = j.Height
                    picPosition = (prs_size[0], Pt(j.Top * 2 + j.Height).emu)
                    if not picNeedDel:
                        picNeedDel = j
                if j.TextFrame.HasText:
                    if "SubTitle" in j.TextFrame.TextRange.Text:
                        j.TextFrame.TextRange.Text = catalogue
                    elif "Introduction" in j.TextFrame.TextRange.Text:
                        introduction = j
                        j.TextFrame.TextRange.Text = image_jianjie
                    elif j.TextFrame.TextRange.ParagraphFormat.Bullet.Visible:
                        j.TextFrame.TextRange.Text = explain
                        Explain = j
            raise RuntimeError("break")
        except RuntimeError:
            break
        except Exception as e:
            print(f"set_ppt_image_explain error {e = }, Try again!")
            picHeight = 100
            introduction = None
            picNeedDel = None
            picPosition = ()
            Explain = None
            try:
                target.Close()
            except Exception:
                pass
            powerpoint.Quit()
            powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
            target = powerpoint.Presentations.Open(target_file)
    if picNeedDel:
        picNeedDel.Delete()
    # picPath = [nf_file[:-5] + "\\ppt\\media\\" + images[i] for i in range(len(images))]
    # picPosition = get_picture_position(picPath, prs_size, picHeight, log)
    # print(f"{picPosition = }")
    # for i in picPosition:
    #     target.Slides(index).Shapes.AddPicture(i[0], True, True, i[1], i[2], i[3], i[4])
    picPath = nf_file[:-5] + "\\ppt\\media\\" + images[0]
    picPos = get_picture_position(picPath, picPosition, picHeight, log)
    print(f"{picPos = }")
    target.Slides(index).Shapes.AddPicture(
        picPos[0], True, True, picPos[1], picPos[2], picPos[3], picPos[4]
    )
    introduction.Left = (Emu(prs_size[0]).pt - introduction.Width) / 2.0000
    Explain.Width = min(Explain.Width, Emu(prs_size[0]).pt - 100)
    Explain.Left = (Emu(prs_size[0]).pt - Explain.Width) / 2.0000
    Explain.Top = introduction.Top + introduction.Height
    target.Save()
    print("SaveTarget!")
    target.Close()
    print("CloseTarget!")
    return


def get_picture_position(
    path, prs_size, picHeight, log
) -> list[tuple[str, int, int, int, int]]:
    """return tuple(path, left, top, width, height)"""
    prs_width, prs_height = prs_size
    prs_width = Emu(prs_width).pt
    prs_height = Emu(prs_height).pt
    with Image.open(path) as img:
        width, height = img.size
    width1 = width / height * picHeight
    height1 = picHeight
    if width1 > prs_width:
        height1 = height / width * (prs_width - 100)
        width1 = prs_width - 100
    return (
        path,
        (prs_width - width1) / 2.0000,
        (prs_height - height1) / 2.0000,
        width1,
        height1,
    )


def get_ppt_only_image_slideIndex(mb_file, log):
    powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
    powerpoint.Visible = 1
    prs = Presentation(mb_file)
    prses = powerpoint.Presentations.Open(mb_file)
    flag = 0
    for i in range(len(prses.Slides)):
        flag = 0
        for j in prses.Slides(i + 1).Shapes:
            if j.Type == MSO_SHAPE_TYPE.PICTURE:
                flag += 1
        if flag > 1:
            for j in prses.Slides(i + 1).Shapes:
                if (
                    j.HasTextFrame
                    and j.TextFrame.TextRange.ParagraphFormat.Bullet.Visible
                ):
                    flag = 0
                    break
        if flag:
            prses.Close()
            del prs
            return i + 1
    prses.Close()
    del prs
    log("can't find only image slide")
    return -1


def get_ppt_only_paragraph_slideIndex(mb_file, log):
    powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
    powerpoint.Visible = 1
    prs = Presentation(mb_file)
    prses = powerpoint.Presentations.Open(mb_file)
    for i in range(len(prses.Slides)):
        flag = 0
        for j in prses.Slides(i + 1).Shapes:
            if j.Type == MSO_SHAPE_TYPE.PICTURE:
                flag += 1
        if flag == 1:
            for j in prses.Slides(i + 1).Shapes:
                if (
                    j.HasTextFrame
                    and j.TextFrame.TextRange.ParagraphFormat.Bullet.Visible
                ):
                    flag = 0
                    break
        if not flag:
            prses.Close()
            del prs
            return i + 1
    prses.Close()
    del prs
    log("can't find only paragraph slide")
    return -1


def get_ppt_image_paragraph_slideIndex(mb_file, log):
    powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
    powerpoint.Visible = 1
    prs = Presentation(mb_file)
    prses = powerpoint.Presentations.Open(mb_file)
    for i in range(len(prses.Slides)):
        flag = 0
        for j in prses.Slides(i + 1).Shapes:
            if j.Type == MSO_SHAPE_TYPE.PICTURE:
                flag += 1
        if flag >= 2:
            for j in prses.Slides(i + 1).Shapes:
                if (
                    j.HasTextFrame
                    and j.TextFrame.TextRange.ParagraphFormat.Bullet.Visible
                ):
                    flag = 0
                    break
        if not flag:
            prses.Close()
            del prs
            return i + 1
    prses.Close()
    del prs
    log("can't find image and paragraph slide")
    return -1


def set_ppt_catalogue(mb_file, target_file, mbIndex, subtitle, catalogues, log):
    time.sleep(1)
    print(f"{subtitle = } {catalogues = }")
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1
    while 10:
        try:
            mb = powerpoint.Presentations.Open(mb_file)
            mb.Slides(mbIndex).Copy()
            mb.Close()
            index = get_needFormatPpt_count(target_file, log) + 1
            target = powerpoint.Presentations.Open(target_file)
            target.Slides.Paste(index)
            for j in target.Slides(index).Shapes:
                if j.TextFrame.HasText:
                    if j.Type == MSO_SHAPE_TYPE.GROUP:
                        j.Ungroup()
                    tmpText = j.TextFrame.TextRange.Text
                    if "提纲" in tmpText or "catalogue" in tmpText:
                        j.TextFrame.TextRange.Text = subtitle
                        continue
                    if j.TextFrame.TextRange.ParagraphFormat.Bullet.Visible:
                        j.TextFrame.TextRange.ParagraphFormat.Bullet.Visible = False
                        j.TextFrame.TextRange.Text = catalogues
            target.Save()
            target.Close()
            raise RuntimeError("break")
        except RuntimeError:
            return
        except Exception as e:
            log(f"set_ppt_catalogue {catalogues} error {e = }")
            print(f"set_ppt_catalogue {catalogues} error {e = }")
            powerpoint.Quit()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")


def get_ppt_catalogues_slideIndex(mb_file, log):
    powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
    powerpoint.Visible = 1
    mb = powerpoint.Presentations.Open(mb_file)
    for i in range(len(mb.Slides)):
        for j in range(len(mb.Slides(i + 1).Shapes)):
            if mb.Slides(i + 1).Shapes(j + 1).Type == MSO_SHAPE_TYPE.GROUP:
                mb.Slides(i + 1).Shapes(j + 1).Ungroup()
            if mb.Slides(i + 1).Shapes(j + 1).TextFrame.HasText:
                text = mb.Slides(i + 1).Shapes(j + 1).TextFrame.TextRange.Text
                if "提纲" in text or "catalogue" in text:
                    mb.Close()
                    powerpoint.Quit()
                    return i + 1
    mb.Close()
    log("can't find catalogue slide")
    powerpoint.Quit()
    return -1


def set_ppt_chapter(mb_file, _, target_file, mbIndex, chapter, log):
    time.sleep(1)
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1
    mb = powerpoint.Presentations.Open(mb_file)
    target = powerpoint.Presentations.Open(target_file)
    mb.Slides(mbIndex).Copy()
    index = get_needFormatPpt_count(target_file, log) + 1
    target.Slides.Paste(index)
    mb.Close()
    for j in range(len(target.Slides(index).Shapes)):
        if target.Slides(index).Shapes(j + 1).TextFrame.HasText:
            if target.Slides(index).Shapes(j + 1).Type == MSO_SHAPE_TYPE.GROUP:
                target.Slides(index).Shapes(j + 1).Ungroup()
            tmpText = target.Slides(index).Shapes(j + 1).TextFrame.TextRange.Text
            if "第" in tmpText and "章" in tmpText:
                target.Slides(index).Shapes(j + 1).TextFrame.TextRange.Text = chapter
                break
    target.Save()
    target.Close()
    return


def get_needFormatPpt_count(need_format_file, log) -> int:
    if need_format_file[-4:] == ".ppt":
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        powerpoint.Visible = 1
        prs = powerpoint.Presentations.Open(need_format_file)
        i = len(prs.Slides)
        prs.Close()
        return i
    else:
        prs = Presentation(need_format_file)
        return len(prs.slides)


def get_ppt_chapter_slideIndex(mb_file, log) -> int:
    powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
    powerpoint.Visible = 1
    mb = powerpoint.Presentations.Open(mb_file)
    for i in range(len(mb.Slides)):
        for j in range(len(mb.Slides(i + 1).Shapes)):
            if mb.Slides(i + 1).Shapes(j + 1).Type == MSO_SHAPE_TYPE.GROUP:
                mb.Slides(i + 1).Shapes(j + 1).Ungroup()
            if mb.Slides(i + 1).Shapes(j + 1).TextFrame.HasText:
                text = mb.Slides(i + 1).Shapes(j + 1).TextFrame.TextRange.Text
                if ("第" in text and "章" in text) or "hapter" in text:  # [C/c]hapter
                    mb.Close()
                    return i + 1
    mb.Close()
    log("can't find chapter slide")
    return -1


def set_ppt_title(mb_file, ppt_file, target_file, log):
    """return tuple(prs.slide_width, prs.slide_height)"""
    powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
    powerpoint.Visible = 1
    prs = Presentation(mb_file)
    title = {"Index": 0, "Font": None, "Size": None, "Color": None}
    for i in range(len(prs.slides)):
        for shape in prs.slides[i].shapes:
            if shape.has_text_frame and "Title" in shape.text:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "Title" in run.text:
                            title["Index"] = i
                            title["Font"] = run.font.name
                            title["Size"] = run.font.size
                            title["Color"] = (
                                run.font.color.rgb
                                if isinstance(run.font.color, _NoneColor)
                                else None
                            )
                            break
    log(f"get ppt_title success\n{title = }")
    tmp = [
        prs.slide_height,
        prs.slide_width,
        prs.slides[title["Index"]],
        prs.slide_layouts[title["Index"]],
    ]
    if not os.path.exists(os.path.dirname(target_file) + "\\"):
        os.mkdir(os.path.dirname(target_file) + "\\")
    target = powerpoint.Presentations.Add()
    target.SaveAs(target_file)
    target.Close()
    target = Presentation(target_file)
    target.slide_height = tmp[0]
    target.slide_width = tmp[1]
    target.save(target_file)
    target = powerpoint.Presentations.Open(target_file)
    prs2 = Presentation(ppt_file)
    titleText = prs2.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
    mb = powerpoint.Presentations.Open(mb_file)
    mb.Slides(1).Copy()
    target.Slides.Paste(1)
    for i in range(len(target.Slides)):
        for j in range(len(target.Slides(i + 1).Shapes)):
            if (
                target.Slides(i + 1).Shapes(j + 1).TextFrame.HasText
                and "Title"
                in target.Slides(i + 1).Shapes(j + 1).TextFrame.TextRange.Text
            ):
                target.Slides(i + 1).Shapes(
                    j + 1
                ).TextFrame.TextRange.Text = titleText.text
    mb.Close()
    target.Save()
    target.Close()
    log(f"save ppt_file success\n{ppt_file = }")
    log("success to format title")
    powerpoint.Quit()
    return tuple([tmp[1], tmp[0]])


def is_number(s) -> bool:
    try:
        float(s)
        return True
    except ValueError:
        pass
    try:
        import unicodedata

        for i in s:
            unicodedata.numeric(i)
            return True
    except (TypeError, ValueError):
        pass
    return False
